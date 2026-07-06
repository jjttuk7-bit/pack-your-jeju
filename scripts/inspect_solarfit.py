"""SolarFit 발표.pptx를 파싱해 슬라이드 구조·톤을 개괄한다.

Windows 콘솔 CP949 문제 회피 위해 파일로 저장 (UTF-8).
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

SRC = Path(r"C:/Users/USER/Downloads/SolarFit_발표.pptx.pptx")
OUT = Path("data/solarfit_outline.txt")
OUT.parent.mkdir(parents=True, exist_ok=True)

p = Presentation(str(SRC))
lines: list[str] = []
lines.append(f"slides: {len(p.slides)}")
lines.append(
    f"slide_size: {p.slide_width}x{p.slide_height} EMU "
    f"({Emu(p.slide_width).inches}x{Emu(p.slide_height).inches} inches)"
)

for i, s in enumerate(p.slides, 1):
    layout = s.slide_layout.name if s.slide_layout else "?"
    lines.append(f"\n--- slide {i} (layout={layout}) ---")

    for shp in s.shapes:
        try:
            w = Emu(shp.width).inches if shp.width else 0
            h = Emu(shp.height).inches if shp.height else 0
            l = Emu(shp.left).inches if shp.left else 0
            t = Emu(shp.top).inches if shp.top else 0
        except Exception:
            w = h = l = t = 0

        if shp.has_text_frame and shp.text_frame.text.strip():
            for para in shp.text_frame.paragraphs:
                text = "".join(r.text for r in para.runs).strip()
                if not text:
                    continue
                fs = None
                fb = None
                fc = None
                if para.runs:
                    r0 = para.runs[0]
                    fs = r0.font.size.pt if r0.font.size else None
                    fb = r0.font.bold
                    try:
                        if r0.font.color and r0.font.color.rgb:
                            fc = str(r0.font.color.rgb)
                    except Exception:
                        pass
                lines.append(
                    f"  TEXT pos=({l:.1f},{t:.1f}) size=({w:.1f}x{h:.1f}) "
                    f"font_size={fs} bold={fb} color={fc}: {text[:200]}"
                )
        else:
            lines.append(
                f"  SHAPE type={shp.shape_type} name={shp.name} "
                f"pos=({l:.1f},{t:.1f}) size=({w:.1f}x{h:.1f})"
            )

OUT.write_text("\n".join(lines), encoding="utf-8")
print(f"written {OUT}  bytes={OUT.stat().st_size}")
