"""Pack Your Jeju 발표 pptx 조립.

참조: 사용자가 첨부한 SolarFit 발표 자료(20 x 11.25 inches, 카드 그리드,
     좌측 상단 로고 · 강조 컬러 통일 · 마지막 회고 슬라이드).

우리 스코프: 3막 10분 (오프닝 1.5 + 본편 5 + 클로징 1.5 + Q&A 2)에 맞춰 18장.
컬러 톤: Pack Your Jeju 프론트와 일치 (Citrus 오렌지 · Basalt · Ivory · Mint).
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# ─── 팔레트 ─────────────────────────────

CITRUS      = RGBColor(0xE6, 0x7A, 0x34)
CITRUS_2    = RGBColor(0xC5, 0x5A, 0x18)
CITRUS_SOFT = RGBColor(0xFF, 0xE7, 0xD3)
MINT        = RGBColor(0x3C, 0x8D, 0x6F)
MINT_SOFT   = RGBColor(0xDD, 0xF0, 0xE6)
AMBER       = RGBColor(0xD9, 0x77, 0x06)
AMBER_SOFT  = RGBColor(0xFE, 0xF3, 0xC7)
BASALT      = RGBColor(0x2D, 0x2A, 0x26)
BASALT_2    = RGBColor(0x5A, 0x55, 0x4F)
IVORY       = RGBColor(0xFD, 0xF6, 0xEA)
BG_SAND     = RGBColor(0xFB, 0xF6, 0xEA)
EARTH       = RGBColor(0xE9, 0xDF, 0xCF)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RED_SOFT    = RGBColor(0xFE, 0xE2, 0xE2)
RED_DARK    = RGBColor(0x7F, 0x1D, 0x1D)

FONT = "맑은 고딕"


# ─── 슬라이드 프레임 ────────────────────

SLIDE_W = Inches(20)
SLIDE_H = Inches(11.25)


def _new_prs() -> Presentation:
    p = Presentation()
    p.slide_width = SLIDE_W
    p.slide_height = SLIDE_H
    return p


def _bg(slide, rgb=BG_SAND):
    """전체 배경 사각형."""
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    r.fill.solid()
    r.fill.fore_color.rgb = rgb
    r.line.fill.background()
    return r


def _txt(slide, text, *, left, top, width, height,
         size=14, bold=False, color=BASALT, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def _multi_txt(slide, lines, *, left, top, width, height,
               size=14, bold=False, color=BASALT, align=PP_ALIGN.LEFT,
               anchor=MSO_ANCHOR.TOP, line_spacing=1.25, bullet=False):
    """여러 줄 (각 줄: 문자열 또는 dict{text,size,bold,color,align})."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        if isinstance(line, dict):
            r.text = ("· " if bullet else "") + line.get("text", "")
            r.font.name = FONT
            r.font.size = Pt(line.get("size", size))
            r.font.bold = line.get("bold", bold)
            c = line.get("color", color)
            r.font.color.rgb = c if isinstance(c, RGBColor) else color
            if line.get("align"):
                p.alignment = line["align"]
        else:
            r.text = ("· " if bullet else "") + str(line)
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
    return tb


def _card(slide, *, left, top, width, height,
          fill=WHITE, line=EARTH, line_w=0.75, radius=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_w)
    # 라운드 정도 조정
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    return shp


def _pill(slide, text, *, left, top, width, height,
          fill=MINT_SOFT, fg=MINT, size=11, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = fg
    shp.line.width = Pt(0.5)
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    tf = shp.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = fg


def _header(slide, title, kicker=None):
    """모든 본문 슬라이드 공통 상단 헤더 — 좌측 상단 브랜드 + 큰 타이틀."""
    # 좌측 브랜드
    _txt(slide, "PACK YOUR JEJU", left=Inches(0.9), top=Inches(0.55),
         width=Inches(6), height=Inches(0.35), size=11, bold=True, color=CITRUS_2)
    # 상단 라인
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.9), Inches(0.95),
                                  Inches(18.2), Emu(20000))
    line.fill.solid()
    line.fill.fore_color.rgb = EARTH
    line.line.fill.background()
    # 큰 타이틀
    if kicker:
        _txt(slide, kicker, left=Inches(0.9), top=Inches(1.2),
             width=Inches(18), height=Inches(0.4), size=13, bold=True, color=MINT)
        _txt(slide, title, left=Inches(0.9), top=Inches(1.55),
             width=Inches(18), height=Inches(0.9), size=32, bold=True, color=BASALT)
    else:
        _txt(slide, title, left=Inches(0.9), top=Inches(1.25),
             width=Inches(18), height=Inches(0.9), size=32, bold=True, color=BASALT)


def _footer(slide, page: int, total: int):
    _txt(slide, "Pack Your Jeju · 아이펠톤 발표",
         left=Inches(0.9), top=Inches(10.7), width=Inches(10), height=Inches(0.3),
         size=10, color=BASALT_2)
    _txt(slide, f"{page:02d} / {total:02d}",
         left=Inches(18.3), top=Inches(10.7), width=Inches(1), height=Inches(0.3),
         size=10, color=BASALT_2, align=PP_ALIGN.RIGHT)


# ─── 슬라이드 빌더 ─────────────────────

TOTAL = 18


def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _bg(s, BASALT)
    # 오른쪽 3분의 1에 대비 컬러 대각선 블록
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(14), 0, Inches(6), SLIDE_H)
    accent.fill.solid()
    accent.fill.fore_color.rgb = CITRUS
    accent.line.fill.background()

    # 좌측 콘텐츠
    _txt(s, "AIFFELTHON 2026", left=Inches(1.2), top=Inches(1.0),
         width=Inches(10), height=Inches(0.5), size=14, bold=True, color=CITRUS)
    _txt(s, "Pack Your Jeju",
         left=Inches(1.2), top=Inches(3.4),
         width=Inches(14), height=Inches(1.8),
         size=88, bold=True, color=WHITE)
    _txt(s, "제주 특화 · 정직한 여행 준비 서비스",
         left=Inches(1.2), top=Inches(5.6),
         width=Inches(14), height=Inches(0.7),
         size=22, color=IVORY)

    # 한 문장 원칙
    _txt(s, "근거 없이 답하지 않는다.",
         left=Inches(1.2), top=Inches(7.0),
         width=Inches(14), height=Inches(0.6),
         size=20, bold=True, color=CITRUS_SOFT)
    _txt(s, "억지로 채우는 것보다 정직한 '확인 불가'가 낫다.",
         left=Inches(1.2), top=Inches(7.6),
         width=Inches(14), height=Inches(0.6),
         size=15, color=IVORY)

    # 팀 정보 (액센트 블록 위)
    _multi_txt(s, [
        {"text": "TEAM", "size": 11, "bold": True, "color": IVORY},
        {"text": "Pack Your Jeju", "size": 20, "bold": True, "color": WHITE},
    ], left=Inches(14.5), top=Inches(9.0), width=Inches(5), height=Inches(1.5),
       line_spacing=1.3)

    # 하단 마감
    _txt(s, "pack-your-jeju.vercel.app  ·  10분 3막 발표",
         left=Inches(1.2), top=Inches(10.5),
         width=Inches(14), height=Inches(0.4),
         size=11, color=IVORY)


def slide_toc(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "10분 3막", kicker="목차")

    items_left = [
        ("01", "훅 · 여행 정보가 늙는 속도"),
        ("02", "1,686건의 증거"),
        ("03", "그래서 저희는 신뢰 엔진을 심었습니다"),
        ("04", "시스템 아키텍처"),
        ("05", "폼 → 팩 · 실제 데모"),
    ]
    items_right = [
        ("06", "정직함 장면 · 감귤 × 7월"),
        ("07", "정직함 장면 · 우도 × 카페"),
        ("08", "하루방 능동 인사 · 여행 저널"),
        ("09", "킥1 /verify · AI가 AI를 팩트체크"),
        ("10", "게이트 · 회고 · 로드맵 · 감사"),
    ]

    def _col(items, x):
        y = Inches(3.0)
        for no, txt in items:
            _txt(s, no, left=x, top=y, width=Inches(1.2), height=Inches(0.6),
                 size=24, bold=True, color=CITRUS)
            _txt(s, txt, left=x + Inches(1.4), top=y + Inches(0.05),
                 width=Inches(8), height=Inches(0.6),
                 size=18, color=BASALT)
            y += Inches(1.0)

    _col(items_left, Inches(0.9))
    _col(items_right, Inches(10.5))

    _footer(s, 2, TOTAL)


def slide_hook(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, BASALT)

    _txt(s, "오프닝 · 훅", left=Inches(0.9), top=Inches(0.55),
         width=Inches(6), height=Inches(0.35), size=11, bold=True, color=CITRUS)

    _txt(s, "여행 앱들은 다 비슷하게 말합니다.",
         left=Inches(1.2), top=Inches(3.0),
         width=Inches(17), height=Inches(1.0),
         size=32, color=IVORY)

    _multi_txt(s, [
        {"text": "\"가성비 좋은 맛집\"", "size": 44, "bold": True, "color": WHITE},
        {"text": "\"숨은 명소\"", "size": 44, "bold": True, "color": WHITE},
        {"text": "\"꼭 가봐야 할 곳\"", "size": 44, "bold": True, "color": WHITE},
    ], left=Inches(1.2), top=Inches(4.3), width=Inches(17), height=Inches(3.6),
       line_spacing=1.15)

    _txt(s, "그중 폐업한 곳이 얼마나 될까요?",
         left=Inches(1.2), top=Inches(8.4),
         width=Inches(17), height=Inches(0.8),
         size=28, bold=True, color=CITRUS)

    _footer(s, 3, TOTAL)


def slide_evidence(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "1,686건 — 이용자가 직접 접수한 정보 정정 요청",
            kicker="실증 데이터 · 킥3")

    # 큰 숫자 세 개
    xs = [Inches(0.9), Inches(7.0), Inches(13.1)]
    triples = [
        ("1,686", "총 수정요청 건수", "비짓제주 (제주 공식 관광 포털)"),
        ("372", "물리적 변화", "폐업 20 · 이전 26 · 시간 변경 235 · 주소 91"),
        ("35%", "top 3만 590건", "운영시간 235 · 상세정보 202 · 메뉴 154"),
    ]
    for x, (big, sub, foot) in zip(xs, triples):
        _card(s, left=x, top=Inches(2.9), width=Inches(5.9), height=Inches(3.5))
        _txt(s, big, left=x + Inches(0.4), top=Inches(3.0), width=Inches(5), height=Inches(1.8),
             size=68, bold=True, color=CITRUS_2)
        _txt(s, sub, left=x + Inches(0.4), top=Inches(4.8), width=Inches(5), height=Inches(0.6),
             size=18, bold=True, color=BASALT)
        _txt(s, foot, left=x + Inches(0.4), top=Inches(5.4), width=Inches(5), height=Inches(0.6),
             size=11, color=BASALT_2)

    # 유형 분포 표
    _txt(s, "유형 분포", left=Inches(0.9), top=Inches(6.9), width=Inches(4), height=Inches(0.5),
         size=14, bold=True, color=MINT)

    rows = [
        ("폐업/영업종료", "20", "1.2%"),
        ("이전/위치변경", "26", "1.5%"),
        ("운영시간/휴무", "235", "13.9%"),
        ("주소/도로명", "91", "5.4%"),
        ("메뉴/가격", "154", "9.1%"),
        ("상세정보", "202", "12.0%"),
        ("사진/이미지", "108", "6.4%"),
        ("연락처", "119", "7.1%"),
    ]
    tbl_top = Inches(7.4)
    row_h = Inches(0.36)
    col_w = [Inches(4.5), Inches(1.5), Inches(1.5)]
    x0 = Inches(0.9)
    for i, (a, b, c) in enumerate(rows):
        y = tbl_top + i * row_h
        # 배경
        row_fill = IVORY if i % 2 == 0 else WHITE
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, y,
                               col_w[0] + col_w[1] + col_w[2], row_h)
        r.fill.solid(); r.fill.fore_color.rgb = row_fill; r.line.fill.background()
        _txt(s, a, left=x0 + Inches(0.15), top=y + Inches(0.03),
             width=col_w[0], height=row_h, size=12, color=BASALT)
        _txt(s, b, left=x0 + col_w[0], top=y + Inches(0.03),
             width=col_w[1], height=row_h, size=12, bold=True, color=CITRUS_2, align=PP_ALIGN.RIGHT)
        _txt(s, c, left=x0 + col_w[0] + col_w[1], top=y + Inches(0.03),
             width=col_w[2], height=row_h, size=12, color=BASALT_2, align=PP_ALIGN.RIGHT)

    # 오른쪽 인용
    quote_box = _card(s, left=Inches(9.5), top=Inches(7.4),
                      width=Inches(9.6), height=Inches(2.5), fill=IVORY, line=EARTH)
    _multi_txt(s, [
        {"text": "\" 공식 데이터조차 이 속도로 낡습니다.", "size": 20, "bold": True, "color": BASALT},
        {"text": "개인 블로그는 어떨까요? \"", "size": 20, "bold": True, "color": BASALT},
        "",
        {"text": "→ 그래서 저희는 짐 싸기 앱에 신뢰 엔진을 심었습니다.", "size": 15, "bold": True, "color": CITRUS_2},
    ], left=Inches(10.0), top=Inches(7.7), width=Inches(9.0), height=Inches(2.0),
       line_spacing=1.35)

    _footer(s, 4, TOTAL)


def slide_promise(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "저희의 세 문장", kicker="Pack Your Jeju 원칙")

    xs = [Inches(0.9), Inches(7.0), Inches(13.1)]
    kickers = ["01", "02", "03"]
    titles = ["출처를 붙인다", "최신만 답한다", "안 지어낸다"]
    bodies = [
        "모든 응답은 공공데이터 근거와 함께 나온다. 사용자는 '왜 이 곳인지' 링크로 확인 가능하다.",
        "valid_until을 스키마 레벨로 강제한다. 못 채우면 DB에 못 들어간다. 감귤이 그 증거.",
        "LLM은 사실을 생성할 권한이 없다. DB가 사실을 대고, LLM은 조립만 한다.",
    ]
    for x, k, t, b in zip(xs, kickers, titles, bodies):
        _card(s, left=x, top=Inches(3.3), width=Inches(5.9), height=Inches(5.5),
              fill=WHITE, line=EARTH)
        _txt(s, k, left=x + Inches(0.5), top=Inches(3.6),
             width=Inches(2), height=Inches(0.8), size=42, bold=True, color=CITRUS)
        _txt(s, t, left=x + Inches(0.5), top=Inches(4.7),
             width=Inches(5), height=Inches(1.0), size=28, bold=True, color=BASALT)
        _txt(s, b, left=x + Inches(0.5), top=Inches(5.9),
             width=Inches(5), height=Inches(2.5), size=14, color=BASALT_2, line_spacing=1.4)

    _txt(s, "발표 내내 반복될 세 단어",
         left=Inches(0.9), top=Inches(9.4), width=Inches(18), height=Inches(0.6),
         size=15, bold=True, color=CITRUS_2)

    _footer(s, 5, TOTAL)


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "결정론 파이프라인 · LLM은 조립만", kicker="시스템 아키텍처")

    # 파이프라인 세로 5단계 (좌)
    steps = [
        ("filters.py", "폼 입력 → 검색 필터", "LLM 없음, 매핑 테이블"),
        ("search.py", "strict → relaxed 완화 재시도", "DB 인덱스 · 하버사인 거리"),
        ("trust.py", "fallback 4분기 + 배지 4종", "인식론 규칙 준수"),
        ("assemble.py", "팩 조립 · 요일별 배치", "LLM 문구 (없으면 템플릿)"),
        ("logging.py", "query_log 적재", "/admin/metrics 라이브"),
    ]
    x = Inches(0.9)
    y = Inches(3.0)
    for i, (name, desc, meta) in enumerate(steps):
        card_h = Inches(1.2)
        _card(s, left=x, top=y, width=Inches(9), height=card_h)
        # 좌측 번호
        _txt(s, f"{i+1:02d}", left=x + Inches(0.3), top=y + Inches(0.25),
             width=Inches(1.2), height=Inches(0.8),
             size=28, bold=True, color=CITRUS)
        _txt(s, name, left=x + Inches(1.5), top=y + Inches(0.2),
             width=Inches(3.5), height=Inches(0.5),
             size=18, bold=True, color=BASALT)
        _txt(s, desc, left=x + Inches(1.5), top=y + Inches(0.65),
             width=Inches(6.8), height=Inches(0.35),
             size=13, color=BASALT_2)
        _txt(s, meta, left=x + Inches(5.2), top=y + Inches(0.2),
             width=Inches(3.6), height=Inches(0.4),
             size=11, color=MINT, align=PP_ALIGN.RIGHT)
        # 화살표 (마지막 제외)
        if i < len(steps) - 1:
            arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                       x + Inches(4.3), y + card_h + Inches(0.05),
                                       Inches(0.4), Inches(0.15))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = EARTH
            arrow.line.fill.background()
        y += card_h + Inches(0.2)

    # 오른쪽 강조 박스: 게이트
    gx = Inches(10.5); gy = Inches(3.0); gw = Inches(8.7); gh = Inches(7.0)
    _card(s, left=gx, top=gy, width=gw, height=gh, fill=BASALT, line=BASALT)
    _txt(s, "GATE", left=gx + Inches(0.5), top=gy + Inches(0.4),
         width=Inches(6), height=Inches(0.5),
         size=13, bold=True, color=CITRUS)
    _txt(s, "packages/eval/run.py", left=gx + Inches(0.5), top=gy + Inches(0.9),
         width=Inches(7), height=Inches(0.7),
         size=26, bold=True, color=WHITE)
    _multi_txt(s, [
        {"text": "verified_precision  ≥  0.9", "size": 18, "color": IVORY, "bold": True},
        {"text": "fallback_accuracy   ≥  0.9", "size": 18, "color": IVORY, "bold": True},
        {"text": "badge_accuracy      ≥  0.8", "size": 18, "color": IVORY, "bold": True},
    ], left=gx + Inches(0.5), top=gy + Inches(2.1), width=Inches(7.7), height=Inches(2.0),
       line_spacing=1.6)

    _txt(s, "미달 시 exit(1) — 배포가 물리적으로 차단됩니다.",
         left=gx + Inches(0.5), top=gy + Inches(4.4),
         width=Inches(7.7), height=Inches(0.8), size=18, bold=True, color=CITRUS_SOFT)

    _txt(s, "최근 실행", left=gx + Inches(0.5), top=gy + Inches(5.4),
         width=Inches(3), height=Inches(0.4), size=12, color=IVORY)
    _txt(s, "12 / 12  ·  3지표 모두 1.00",
         left=gx + Inches(0.5), top=gy + Inches(5.8),
         width=Inches(7.7), height=Inches(0.8), size=28, bold=True, color=MINT_SOFT)

    _footer(s, 6, TOTAL)


def slide_form_flow(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "폼 → 근거 있는 팩 → 요일별 일정", kicker="본편 · 실사용 흐름")

    # 3화면 카드 나열
    cards = [
        ("화면 1", "여행 계획하기",
         ["지역 다중 선택 (12개 값)",
          "출발일 · 기간 · 동행자",
          "목적 · 특별 노트"]),
        ("화면 2", "순간 고르기",
         ["오름 · 바다 · 노을 · 시장",
          "맛집 · 카페 · 곶자왈 · 감귤",
          "카드 다중 선택"]),
        ("화면 3", "팩 결과 (요일별)",
         ["Day 1~N × 지역별 배치",
          "배지 4종 · 근거 URL",
          "카드 클릭 → 상세 확장"]),
    ]
    x_starts = [Inches(0.9), Inches(7.0), Inches(13.1)]
    for x, (k, t, body) in zip(x_starts, cards):
        _card(s, left=x, top=Inches(3.0), width=Inches(5.9), height=Inches(5.8),
              fill=WHITE, line=EARTH)
        _pill(s, k, left=x + Inches(0.5), top=Inches(3.2),
              width=Inches(1.5), height=Inches(0.5), fill=CITRUS_SOFT, fg=CITRUS_2)
        _txt(s, t, left=x + Inches(0.5), top=Inches(3.9),
             width=Inches(5), height=Inches(0.7),
             size=22, bold=True, color=BASALT)
        _multi_txt(s, body, left=x + Inches(0.5), top=Inches(4.8),
                   width=Inches(5), height=Inches(3.5),
                   size=13, color=BASALT_2, line_spacing=1.5, bullet=True)

    # 하단 강조
    _txt(s, "\"각 장소에 배지가 붙습니다. 초록 = 확인됨, 노랑 = 이용자 수정요청 이력. "
         "각 카드에 원본 데이터 링크가 있습니다.\"",
         left=Inches(0.9), top=Inches(9.4), width=Inches(18), height=Inches(0.8),
         size=15, color=CITRUS_2)

    _footer(s, 7, TOTAL)


def slide_honesty_citrus(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "감귤 × 7월", kicker="정직함 장면 ①")

    # 좌: 잘못된 답 예시 (일반 GPT 톤)
    _card(s, left=Inches(0.9), top=Inches(3.0),
          width=Inches(9), height=Inches(6.5), fill=RED_SOFT, line=RED_DARK)
    _pill(s, "일반 GPT", left=Inches(1.2), top=Inches(3.3),
          width=Inches(1.8), height=Inches(0.5), fill=RED_DARK, fg=WHITE)
    _txt(s, "\"7월에 감귤 체험 가능한 농장 3곳을 추천해드리겠습니다.\"",
         left=Inches(1.2), top=Inches(4.2), width=Inches(8.4), height=Inches(1.5),
         size=20, bold=True, color=RED_DARK, line_spacing=1.35)
    _multi_txt(s, [
        "· ○○농장 (주소: ...)",
        "· △△체험장 (주소: ...)",
        "· ◇◇마을 (주소: ...)",
    ], left=Inches(1.2), top=Inches(6.0), width=Inches(8.4), height=Inches(2.0),
       size=16, color=BASALT, line_spacing=1.5)
    _txt(s, "감귤 수확기 = 10월 ~ 이듬해 1월. 7월엔 존재하지 않습니다.",
         left=Inches(1.2), top=Inches(8.2), width=Inches(8.4), height=Inches(0.7),
         size=14, bold=True, color=RED_DARK)

    # 우: 우리 응답
    _card(s, left=Inches(10.2), top=Inches(3.0),
          width=Inches(9), height=Inches(6.5), fill=MINT_SOFT, line=MINT)
    _pill(s, "Pack Your Jeju", left=Inches(10.5), top=Inches(3.3),
          width=Inches(2.5), height=Inches(0.5), fill=MINT, fg=WHITE)
    _txt(s, "\"저희는 7월에 감귤 체험을 지어내지 않습니다.\"",
         left=Inches(10.5), top=Inches(4.2), width=Inches(8.4), height=Inches(1.5),
         size=20, bold=True, color=MINT, line_spacing=1.35)
    _multi_txt(s, [
        "seasonal 정보 유형",
        "valid_until을 스키마로 강제",
        "시즌 종료일 못 채우면 place에 들어가지 못함",
        "→ '이번 시즌에 확인된 프로그램이 없습니다'",
    ], left=Inches(10.5), top=Inches(6.0), width=Inches(8.4), height=Inches(3.0),
       size=14, color=BASALT, line_spacing=1.6, bullet=True)

    _footer(s, 8, TOTAL)


def slide_honesty_udo(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "우도 × 조용한 카페", kicker="정직함 장면 ② · coverage_gap")

    # 4분기 표
    _txt(s, "저희는 '모른다'를 네 가지로 구분합니다.",
         left=Inches(0.9), top=Inches(3.0), width=Inches(18), height=Inches(0.7),
         size=20, bold=True, color=CITRUS_2)

    quads = [
        ("out_of_scope", "제주 여행 정보 범위 밖",
         "'제주 여행 정보 범위 밖입니다'", BASALT_2),
        ("contradicted", "반증이 있음 (폐업/이전 확인)",
         "'폐업/변경이 확인됩니다'  ← 유일하게 단언 허용", CITRUS_2),
        ("retrieval_miss", "DB엔 있는데 못 찾음",
         "사용자 노출 전 재시도로 해소 시도", BASALT_2),
        ("coverage_gap", "저희 데이터에 없음",
         "'저희가 참조하는 공공데이터 기준으로 확인되지 않습니다'", CITRUS_2),
    ]
    positions = [
        (Inches(0.9), Inches(4.0)),
        (Inches(10.0), Inches(4.0)),
        (Inches(0.9), Inches(7.0)),
        (Inches(10.0), Inches(7.0)),
    ]
    for (x, y), (key, mean, quote, accent) in zip(positions, quads):
        _card(s, left=x, top=y, width=Inches(9.1), height=Inches(2.8), fill=WHITE, line=EARTH)
        _pill(s, key, left=x + Inches(0.3), top=y + Inches(0.3),
              width=Inches(2.5), height=Inches(0.5),
              fill=CITRUS_SOFT if accent == CITRUS_2 else MINT_SOFT,
              fg=accent)
        _txt(s, mean, left=x + Inches(0.3), top=y + Inches(1.0),
             width=Inches(8.5), height=Inches(0.6),
             size=17, bold=True, color=BASALT)
        _txt(s, quote, left=x + Inches(0.3), top=y + Inches(1.7),
             width=Inches(8.5), height=Inches(0.9),
             size=13, color=BASALT_2, line_spacing=1.4)

    _txt(s, "인식론 규칙 · 우리 DB ≠ 세계 전체 · 'contradicted'만 단언 허용",
         left=Inches(0.9), top=Inches(10.05), width=Inches(18), height=Inches(0.5),
         size=13, bold=True, color=MINT)

    _footer(s, 9, TOTAL)


def slide_haruban(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "하루방 · 폼 옆에 사는 정직 상담사", kicker="본편 · 능동 인사")

    # 좌 (설명)
    _card(s, left=Inches(0.9), top=Inches(3.0), width=Inches(9), height=Inches(7),
          fill=WHITE, line=EARTH)
    _txt(s, "폼에서 지역+순간을 처음 고르는 순간,\n하루방이 스스로 말을 겁니다.",
         left=Inches(1.3), top=Inches(3.4), width=Inches(8.2), height=Inches(2.0),
         size=22, bold=True, color=BASALT, line_spacing=1.4)

    _multi_txt(s, [
        {"text": "POST /agent/intro (결정적 · LLM 사실 생성 금지)", "size": 12, "bold": True, "color": MINT},
        "",
        "1. build_filters(form) → search_strict",
        "2. 각 (region × moment) 조합에서 verified 우선 하이라이트 6개",
        "3. gaps 나열 — items=0 조합 정직 노출",
        "4. LLM은 greeting과 reason 한 줄만",
        "5. candidates 밖 external_id는 무조건 폐기",
    ], left=Inches(1.3), top=Inches(5.6), width=Inches(8.2), height=Inches(4),
       size=13, color=BASALT_2, line_spacing=1.55)

    # 우 (카드 예시)
    _card(s, left=Inches(10.2), top=Inches(3.0), width=Inches(9), height=Inches(7),
          fill=IVORY, line=EARTH)
    _txt(s, "🗿 하루방", left=Inches(10.6), top=Inches(3.3),
         width=Inches(4), height=Inches(0.5), size=15, bold=True, color=BASALT)
    _txt(s, "제주 여행길 지킴이",
         left=Inches(10.6), top=Inches(3.7), width=Inches(6), height=Inches(0.4),
         size=11, color=BASALT_2)

    _card(s, left=Inches(10.6), top=Inches(4.4), width=Inches(8.2), height=Inches(1.4),
          fill=WHITE, line=CITRUS_SOFT)
    _txt(s, "구좌에서 연인과 힐링 여행이시라구요?",
         left=Inches(10.9), top=Inches(4.55), width=Inches(7.6), height=Inches(0.5),
         size=13, bold=True, color=BASALT)
    _txt(s, "저희 공공데이터로 확인된 6곳을 먼저 보여드릴게요.",
         left=Inches(10.9), top=Inches(5.05), width=Inches(7.6), height=Inches(0.5),
         size=12, color=BASALT_2)

    # 하이라이트 카드 예시 두 개
    _card(s, left=Inches(10.6), top=Inches(6.0), width=Inches(8.2), height=Inches(1.7),
          fill=WHITE, line=EARTH)
    _pill(s, "verified", left=Inches(17.4), top=Inches(6.15),
          width=Inches(1.3), height=Inches(0.4), fill=MINT_SOFT, fg=MINT)
    _txt(s, "구좌 · 조용한 카페",
         left=Inches(10.9), top=Inches(6.15), width=Inches(6), height=Inches(0.4),
         size=10, color=BASALT_2)
    _txt(s, "카페달보름", left=Inches(10.9), top=Inches(6.55),
         width=Inches(6), height=Inches(0.5), size=16, bold=True, color=BASALT)
    _txt(s, "제주특별자치도 제주시 구좌읍 대수길 24",
         left=Inches(10.9), top=Inches(7.05), width=Inches(7.5), height=Inches(0.4),
         size=11, color=BASALT_2)

    _card(s, left=Inches(10.6), top=Inches(7.9), width=Inches(8.2), height=Inches(1.7),
          fill=WHITE, line=EARTH)
    _pill(s, "caution", left=Inches(17.4), top=Inches(8.05),
          width=Inches(1.3), height=Inches(0.4), fill=AMBER_SOFT, fg=AMBER)
    _txt(s, "구좌 · 노을 감상",
         left=Inches(10.9), top=Inches(8.05), width=Inches(6), height=Inches(0.4),
         size=10, color=BASALT_2)
    _txt(s, "메이즈랜드", left=Inches(10.9), top=Inches(8.45),
         width=Inches(6), height=Inches(0.5), size=16, bold=True, color=BASALT)
    _txt(s, "· 이용자 정보 수정요청 이력",
         left=Inches(10.9), top=Inches(8.95), width=Inches(7.5), height=Inches(0.4),
         size=11, color=AMBER)

    _footer(s, 10, TOTAL)


def slide_journal(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "이 여행을, 저널로 저장", kicker="본편 · POST /pack/pdf")

    _multi_txt(s, [
        {"text": "\"근거 있는 것만 담고, 없는 것은 정직하게 비웠습니다.\"", "size": 24, "bold": True, "color": BASALT},
        "",
        {"text": "표지 → Day별 카드 → 마지막 페이지: '정직하게, 확인되지 않은 것들'", "size": 15, "color": BASALT_2},
    ], left=Inches(0.9), top=Inches(3.0), width=Inches(18), height=Inches(2.0),
       line_spacing=1.35)

    boxes = [
        ("표지",
         ["대문 문구: 이 여행을, 정직하게 담았습니다.",
          "부제: 동행자 · 지역 · 기간 · 목적",
          "커버리지 chip: 확인됨 · caution · gap"]),
        ("Day별 카드",
         ["요일 · 지역 · 순간 라벨",
          "이름 · 배지 · 주소 · 근거 URL",
          "미확인 조합 amber 노트"]),
        ("정직 페이지",
         ["'저희 데이터에서 확인되지 않은 조합'",
          "gaps 요약 · '없다'로 단언하지 않음",
          "마무리: '근거 있는 여행 준비'"]),
    ]
    xs = [Inches(0.9), Inches(7.0), Inches(13.1)]
    for x, (t, body) in zip(xs, boxes):
        _card(s, left=x, top=Inches(5.3), width=Inches(5.9), height=Inches(4.7))
        _txt(s, t, left=x + Inches(0.4), top=Inches(5.6),
             width=Inches(5), height=Inches(0.7),
             size=20, bold=True, color=CITRUS_2)
        _multi_txt(s, body, left=x + Inches(0.4), top=Inches(6.6),
                   width=Inches(5.2), height=Inches(3.2),
                   size=13, color=BASALT_2, line_spacing=1.55, bullet=True)

    _footer(s, 11, TOTAL)


def slide_verify(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "AI가 AI를 팩트체크", kicker="킥1 · POST /verify")

    _txt(s, "리뷰를 입력값으로 받아 문장별로 판정 · 근거 URL 반환",
         left=Inches(0.9), top=Inches(3.0), width=Inches(18), height=Inches(0.6),
         size=18, color=BASALT)

    # 리뷰 원문
    _card(s, left=Inches(0.9), top=Inches(4.0), width=Inches(18.2), height=Inches(1.9),
          fill=IVORY, line=EARTH)
    _txt(s, "입력 리뷰 (킥1 대본)", left=Inches(1.2), top=Inches(4.15),
         width=Inches(8), height=Inches(0.4),
         size=11, bold=True, color=MINT)
    _txt(s, "\"애월오누이 제주 정말 맛있어요. 라임오렌지카페앤플라워는 조용해서 좋았어요. 그 옆 가짜식당은 별로였어요.\"",
         left=Inches(1.2), top=Inches(4.55), width=Inches(17.6), height=Inches(1.3),
         size=17, color=BASALT, line_spacing=1.5)

    # 판정 결과 세 카드
    verdicts = [
        ("애월오누이 제주 정말 맛있어요.",
         "×  contradicted",
         "폐업/이전이 확인됩니다",
         RED_SOFT, RED_DARK, "비짓제주 데이터 tombstoned=true"),
        ("라임오렌지카페앤플라워는 조용해서 좋았어요.",
         "○  verified",
         "저희 데이터에서 존재 확인",
         MINT_SOFT, MINT, "place external_id 매칭"),
        ("그 옆 가짜식당은 별로였어요.",
         "⚠  coverage_gap",
         "저희 데이터로 확인되지 않습니다",
         AMBER_SOFT, AMBER, "'없다'로 단언하지 않음"),
    ]
    xs = [Inches(0.9), Inches(7.0), Inches(13.1)]
    for x, (sent, verd, msg, bg, fg, meta) in zip(xs, verdicts):
        _card(s, left=x, top=Inches(6.3), width=Inches(5.9), height=Inches(3.9),
              fill=bg, line=fg)
        _txt(s, verd, left=x + Inches(0.4), top=Inches(6.5),
             width=Inches(5), height=Inches(0.6),
             size=18, bold=True, color=fg)
        _txt(s, sent, left=x + Inches(0.4), top=Inches(7.2),
             width=Inches(5.1), height=Inches(1.4),
             size=13, color=BASALT, line_spacing=1.4)
        _txt(s, msg, left=x + Inches(0.4), top=Inches(8.7),
             width=Inches(5.1), height=Inches(0.6),
             size=13, bold=True, color=fg)
        _txt(s, meta, left=x + Inches(0.4), top=Inches(9.4),
             width=Inches(5.1), height=Inches(0.6),
             size=10, color=BASALT_2)

    _footer(s, 12, TOTAL)


def slide_badges(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "배지 · 정직함의 시각 자산", kicker="트러스트 엔진 · 항목 판정")

    badges = [
        ("verified",  "확인됨",       "존재 + 유효기간 유효 + 수정요청 없음",
         MINT_SOFT, MINT),
        ("caution",   "주의",         "수정요청 이력 or 요구 amenity 결측 or 유효기간 임박",
         AMBER_SOFT, AMBER),
        ("contradicted", "폐업/변경 확인",
         "tombstone=true 또는 폐업/이전 반증",
         RED_SOFT, RED_DARK),
        ("reference", "참고",         "공공데이터 검증 아님 (스트레치, verified와 명확 분리)",
         RGBColor(0xEE, 0xE7, 0xD6), BASALT_2),
    ]
    positions = [
        (Inches(0.9), Inches(3.0)),
        (Inches(10.0), Inches(3.0)),
        (Inches(0.9), Inches(6.8)),
        (Inches(10.0), Inches(6.8)),
    ]
    for (x, y), (key, ko, cond, bg, fg) in zip(positions, badges):
        _card(s, left=x, top=y, width=Inches(9.1), height=Inches(3.5), fill=WHITE, line=EARTH)
        _pill(s, key, left=x + Inches(0.4), top=y + Inches(0.4),
              width=Inches(2.5), height=Inches(0.6), fill=bg, fg=fg)
        _txt(s, ko, left=x + Inches(0.4), top=y + Inches(1.15),
             width=Inches(8.3), height=Inches(0.9),
             size=32, bold=True, color=BASALT)
        _txt(s, cond, left=x + Inches(0.4), top=y + Inches(2.15),
             width=Inches(8.3), height=Inches(1.1),
             size=14, color=BASALT_2, line_spacing=1.5)

    _footer(s, 13, TOTAL)


def slide_data_stack(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "우리가 사실을 대는 방식", kicker="데이터 스택")

    # 큰 숫자 카드 4개
    numbers = [
        ("4,422", "관광지 · 카페 · 음식점", "visitjeju API 4,422건"),
        ("5,828", "교통 배지 근거", "주차장 1,557 + 정류장 4,271"),
        ("556",   "수정요청 매칭",   "1,686건 중 place에 이력 부착"),
        ("12/12", "골든셋 게이트",   "3지표 모두 1.00 (verified/fallback/badge)"),
    ]
    xs = [Inches(0.9), Inches(5.4), Inches(9.9), Inches(14.4)]
    for x, (big, sub, foot) in zip(xs, numbers):
        _card(s, left=x, top=Inches(3.0), width=Inches(4.3), height=Inches(3.5))
        _txt(s, big, left=x + Inches(0.3), top=Inches(3.2), width=Inches(3.8), height=Inches(1.6),
             size=56, bold=True, color=CITRUS_2)
        _txt(s, sub, left=x + Inches(0.3), top=Inches(4.9), width=Inches(3.8), height=Inches(0.6),
             size=15, bold=True, color=BASALT)
        _txt(s, foot, left=x + Inches(0.3), top=Inches(5.55), width=Inches(3.8), height=Inches(0.7),
             size=11, color=BASALT_2)

    # 하단 소스 표
    _txt(s, "소스",
         left=Inches(0.9), top=Inches(7.0), width=Inches(6), height=Inches(0.5),
         size=14, bold=True, color=MINT)
    sources = [
        ("S1", "비짓제주 API (키)",             "관광지·카페·음식·체험",       "place 4,422"),
        ("S2", "위생등급 CSV",                  "local_food 배지 보강",         "food.hygiene_grade"),
        ("S3", "공영주차장 CSV (공공데이터)",   "교통 배지 🅿️",                "transit_point 1,557"),
        ("S4", "TAGO 버스정류소 (국토부)",      "교통 배지 🚌",                 "transit_point 4,271"),
        ("S5", "비짓제주 수정요청 CSV",         "신뢰 하향 신호 + 킥3 근거",     "has_fix_request 556"),
    ]
    y = Inches(7.5)
    for i, (k, name, role, out) in enumerate(sources):
        row_fill = IVORY if i % 2 == 0 else WHITE
        rr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), y,
                                Inches(18.2), Inches(0.5))
        rr.fill.solid(); rr.fill.fore_color.rgb = row_fill; rr.line.fill.background()
        _txt(s, k, left=Inches(1.0), top=y + Inches(0.05),
             width=Inches(0.8), height=Inches(0.4),
             size=12, bold=True, color=CITRUS_2)
        _txt(s, name, left=Inches(1.9), top=y + Inches(0.05),
             width=Inches(6.0), height=Inches(0.4),
             size=12, color=BASALT)
        _txt(s, role, left=Inches(8.0), top=y + Inches(0.05),
             width=Inches(6.0), height=Inches(0.4),
             size=12, color=BASALT_2)
        _txt(s, out, left=Inches(14.1), top=y + Inches(0.05),
             width=Inches(5.0), height=Inches(0.4),
             size=12, bold=True, color=MINT)
        y += Inches(0.5)

    _footer(s, 14, TOTAL)


def slide_live(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "저희는 발표하는 동안에도 운영 중이었습니다",
            kicker="킥4 · GET /admin/metrics")

    # 왼쪽 QR 안내 카드
    _card(s, left=Inches(0.9), top=Inches(3.0), width=Inches(8.5), height=Inches(7),
          fill=BASALT, line=BASALT)
    _txt(s, "직접 만들어보세요", left=Inches(1.3), top=Inches(3.4),
         width=Inches(7.5), height=Inches(0.9), size=32, bold=True, color=WHITE)
    _txt(s, "pack-your-jeju.vercel.app", left=Inches(1.3), top=Inches(4.4),
         width=Inches(7.5), height=Inches(0.7), size=22, bold=True, color=CITRUS_SOFT)
    # QR 자리표시
    qr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            Inches(2.4), Inches(5.3), Inches(5.5), Inches(4.3))
    qr.fill.solid(); qr.fill.fore_color.rgb = WHITE
    qr.line.color.rgb = WHITE
    _txt(s, "docs/qr.png 삽입",
         left=Inches(2.4), top=Inches(7.2), width=Inches(5.5), height=Inches(0.5),
         size=14, color=BASALT_2, align=PP_ALIGN.CENTER)

    # 오른쪽 라이브 지표
    _txt(s, "최근 24시간 실 트래픽",
         left=Inches(10.0), top=Inches(3.0), width=Inches(9), height=Inches(0.6),
         size=15, bold=True, color=MINT)
    metrics = [
        ("25", "/pack 요청"),
        ("8",  "/verify 요청"),
        ("31", "verified"),
        ("102", "caution"),
        ("6",   "coverage_gap"),
        ("2",   "contradicted"),
    ]
    x0 = Inches(10.0); y0 = Inches(3.8)
    for i, (big, sub) in enumerate(metrics):
        cx = x0 + Inches(3.0) * (i % 3)
        cy = y0 + Inches(2.4) * (i // 3)
        _card(s, left=cx, top=cy, width=Inches(2.9), height=Inches(2.2), fill=WHITE, line=EARTH)
        _txt(s, big, left=cx + Inches(0.2), top=cy + Inches(0.3),
             width=Inches(2.5), height=Inches(1.2),
             size=40, bold=True, color=CITRUS_2)
        _txt(s, sub, left=cx + Inches(0.2), top=cy + Inches(1.5),
             width=Inches(2.5), height=Inches(0.6),
             size=12, color=BASALT_2)

    _txt(s, "지연 p50 = 16ms · 배지 5종 실 데이터로 회전 중",
         left=Inches(10.0), top=Inches(9.4), width=Inches(9.1), height=Inches(0.5),
         size=13, bold=True, color=MINT)

    _footer(s, 15, TOTAL)


def slide_roadmap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "지역 팩 아키텍처", kicker="다음 단계 · Roadmap")

    boxes = [
        ("01",
         "지역 확장",
         "부산 팩 · 도쿄 팩",
         "데이터 어댑터만 갈아끼우면 됩니다. 저희는 모든 곳을 아는 척하는 대신 제주부터 완벽하게 하기로 했습니다."),
        ("02",
         "여행 일정 전체 검증",
         "일정 단위 팩트체크",
         "지금은 개별 조각(place·food)을 검증. 다음은 사용자가 만든 하루/전 일정을 통째로 팩트체크."),
        ("03",
         "다국어 · locale 확장",
         "비짓제주 en / cn / jp",
         "우리 스키마는 이미 locale-agnostic. UI 문구만 로컬라이즈하면 해외 여행자용으로 확장 가능."),
    ]
    xs = [Inches(0.9), Inches(7.0), Inches(13.1)]
    for x, (num, t, sub, body) in zip(xs, boxes):
        _card(s, left=x, top=Inches(3.0), width=Inches(5.9), height=Inches(7))
        _txt(s, num, left=x + Inches(0.4), top=Inches(3.3),
             width=Inches(2), height=Inches(1.0), size=44, bold=True, color=CITRUS)
        _txt(s, t, left=x + Inches(0.4), top=Inches(4.6),
             width=Inches(5.2), height=Inches(0.8),
             size=22, bold=True, color=BASALT)
        _txt(s, sub, left=x + Inches(0.4), top=Inches(5.5),
             width=Inches(5.2), height=Inches(0.7),
             size=14, bold=True, color=MINT)
        _txt(s, body, left=x + Inches(0.4), top=Inches(6.4),
             width=Inches(5.2), height=Inches(3.4),
             size=13, color=BASALT_2, line_spacing=1.55)

    _footer(s, 16, TOTAL)


def slide_retro(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "Keep · Problem · Try", kicker="회고 (Retrospective)")

    columns = [
        ("Keep — 잘한 점",
         MINT,
         [
             "1,686건 실증 데이터를 오프닝에 넣은 것 — 문제 정의 강도 확 오름.",
             "LLM에 사실 생성 권한을 안 준 원칙 · 예외 없이 지킨 것.",
             "골든셋 12문항 게이트를 실제 exit(1)로 코딩 · '말이 아니라 화면'.",
             "하루방 캐릭터로 정직 원칙을 사용자에게 전달 · UX가 정체성과 정합.",
             "여행 저널 PDF · 근거를 실물로 남길 수 있게 함.",
         ]),
        ("Problem — 아쉬운 점",
         AMBER,
         [
             "위생등급 CSV 미적재 · 배지 하나가 완성 안 됨 (P2 스킵).",
             "Railway Docker 빌드 실패(fonts-nanum-coding)로 배포 지연 1시간.",
             "OPENAI_API_KEY 미설정 상태로 LLM 감성 문구는 템플릿 폴백만.",
             "커버리지가 아직 제주 12개 지역에 한정 — 확장 로드맵으로 정직 안내.",
         ]),
        ("Try — 다음에 시도할 것",
         CITRUS_2,
         [
             "지역 팩 어댑터 계층 추상화 · 부산 팩 프로토타입.",
             "여행 일정 전체 팩트체크 (개별 조각 → 조합).",
             "OPENAI_API_KEY 확보 후 감성 문구 · 하루방 reason 라이브 활성.",
             "골든셋 문항 확장 (15 → 30) · 실증 케이스 커버.",
         ]),
    ]
    xs = [Inches(0.9), Inches(7.0), Inches(13.1)]
    for x, (t, fg, items) in zip(xs, columns):
        _card(s, left=x, top=Inches(3.0), width=Inches(5.9), height=Inches(7))
        _pill(s, t, left=x + Inches(0.4), top=Inches(3.3),
              width=Inches(5.1), height=Inches(0.6),
              fill=MINT_SOFT if fg == MINT else (AMBER_SOFT if fg == AMBER else CITRUS_SOFT),
              fg=fg, size=13)
        _multi_txt(s, items, left=x + Inches(0.4), top=Inches(4.2),
                   width=Inches(5.1), height=Inches(5.6),
                   size=13, color=BASALT, line_spacing=1.55, bullet=True)

    _footer(s, 17, TOTAL)


def slide_thanks(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, BASALT)
    # 액센트 대각 블록
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                0, Inches(9.5), SLIDE_W, Inches(0.15))
    accent.fill.solid(); accent.fill.fore_color.rgb = CITRUS
    accent.line.fill.background()

    _txt(s, "감사합니다.",
         left=Inches(1.2), top=Inches(3.6),
         width=Inches(17), height=Inches(1.8),
         size=88, bold=True, color=WHITE)

    _multi_txt(s, [
        {"text": "\"근거가 있는 것만 말합니다.", "size": 22, "color": CITRUS_SOFT, "bold": True},
        {"text": " 최신 정보만 답합니다.", "size": 22, "color": CITRUS_SOFT, "bold": True},
        {"text": " 지어내지 않습니다.\"", "size": 22, "color": CITRUS_SOFT, "bold": True},
    ], left=Inches(1.2), top=Inches(5.7), width=Inches(17), height=Inches(2.3),
       line_spacing=1.4)

    _txt(s, "Pack Your Jeju  ·  pack-your-jeju.vercel.app",
         left=Inches(1.2), top=Inches(9.9),
         width=Inches(17), height=Inches(0.6),
         size=16, color=IVORY)


# ─── 조립 ─────────────────────────────

def build() -> Path:
    prs = _new_prs()
    slide_cover(prs)          # 1
    slide_toc(prs)            # 2
    slide_hook(prs)           # 3
    slide_evidence(prs)       # 4
    slide_promise(prs)        # 5
    slide_architecture(prs)   # 6
    slide_form_flow(prs)      # 7
    slide_honesty_citrus(prs) # 8
    slide_honesty_udo(prs)    # 9
    slide_haruban(prs)        # 10
    slide_journal(prs)        # 11
    slide_verify(prs)         # 12
    slide_badges(prs)         # 13
    slide_data_stack(prs)     # 14
    slide_live(prs)           # 15
    slide_roadmap(prs)        # 16
    slide_retro(prs)          # 17
    slide_thanks(prs)         # 18

    out = Path("docs") / "pack_your_jeju_slides.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


if __name__ == "__main__":
    p = build()
    print(f"generated: {p.resolve()}  size={p.stat().st_size:,} bytes  slides={TOTAL}")
